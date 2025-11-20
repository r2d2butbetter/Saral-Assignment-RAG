import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io


def extract_cited_sources(text):
    """Extract cited source numbers from text."""
    return set(re.findall(r'\[Source (\d+)\]', text))


def create_presentation(slide_data):
    """Create a PowerPoint presentation from SlideData object."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = slide_data.presentation_title
    
    # Content slides
    bullet_slide_layout = prs.slide_layouts[1]
    for slide_content in slide_data.slides:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = slide_content.title
        
        text_frame = body_shape.text_frame
        text_frame.clear()
        
        for bullet_text in slide_content.bullets:
            p = text_frame.add_paragraph()
            p.text = bullet_text
            p.level = 0
    
    # Save to BytesIO buffer
    pptx_buffer = io.BytesIO()
    prs.save(pptx_buffer)
    pptx_buffer.seek(0)
    
    return pptx_buffer
